"""

Para aprender a generar PPTx desde Python ir a: https://pbpython.com/creating-powerpoint.html

Tip: también existen herramientas para generar un PPT desde un archivo Markdown

"""
from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches

# template
prs = Presentation("160245-books-template-16x9.pptx")

def set_title(slide, title_text):
    try:
        title = slide.shapes.title
        title.text = title_text
    except AttributeError:
        print("No Title for Layout!")

def add_text(slide, placeholder_index, shape_name, text):
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            # Do not overwrite the title which is just a special placeholder
            try:
                if 'Title' not in shape.text:
                    if phf.idx == placeholder_index and shape.name == shape_name:
                        shape.text = text
            except AttributeError:
                print("{} has no text attribute".format(phf.type))
            print('{} {}'.format(phf.idx, shape.name))

def remove_placeholder(slide, placeholder_index):
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.idx == placeholder_index:
                phf = None
                return

# welcome page
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_title(slide,"Documentation as code\nHacktoberFest 2020")
add_text(slide,1,"Subtitle 2","Henry Tong\nSoftware Architect - Backend Developer\ntaksantong@gmail.com")

# documentation is the last thing....
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide,"Entregable de Software")
add_text(slide,1,"Content Placeholder 2","Siempre: Código fuente (*.java *.js *.py ...)\nA veces: Ejecutable\n¿Alguna vez?: Documentación")

# quote
slide = prs.slides.add_slide(prs.slide_layouts[2])
left=Inches(2)
top=Inches(1)
image_path='fool.png'
img=slide.shapes.add_picture(image_path,left,top)

# documentation as code....
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide,"Documentation as code")
ptext = ("Documentar mientras se programa, usando la IDE/Editor habitual\n"
        "Se prefiere formato de texto (ReStructuredText, ASCIIDoc, Markdown...)\n"
        "Forma parte del proceso de desarrollo: commit, push, deploy...\n"
        "Si se puede automatizar mucho mejor! Por ejemplo esta presentación también es código!")
add_text(slide,1,"Content Placeholder 2",ptext)

# Además!
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide,"También!")
add_text(slide,1,"Content Placeholder 2","Container as code: Dockerfile\nDiagram as code: Diagrams, PlantUML\nInfrastructure as code: Especificar servidores, nodos, .etc. de una arquitectura: Terraform, Cloudformation...")

# Demo
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide,"Demo!")
add_text(slide,1,"Content Placeholder 2","Vamos a codear! (y documentar!)")

# Herramientas 
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide,"Herramientas recomendadas")
tools=("Swagger para documentar APIs https://swagger.io\n"
"Visual Studio Code plugin para probar REST https://marketplace.visualstudio.com/items?itemName=humao.rest-client\n"
"Documentador Docosaurus hecho en javascript https://docusaurus.io\n"
"Documentador MkDocs hecho en python https://www.mkdocs.org\n"
"Más generadores de web estáticas buscando en google \"static site generator\"")
add_text(slide,1,"Content Placeholder 2",tools)

# Gracias!
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide,"Gracias")
add_text(slide,1,"Content Placeholder 2","Para conocer más sobre \"Documentation as code\" https://www.writethedocs.org/guide/docs-as-code")

# save result
prs.save("generated_presentation.pptx")

print("done!")
